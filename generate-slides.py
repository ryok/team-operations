"""
松尾研究所テンプレート準拠のスライド生成スクリプト

使い方:
    uvx --from python-pptx python generate-slides.py

テンプレートファイルを開き、既存サンプルスライドを削除してから
新しいスライドを追加する。スライドマスターのロゴ・コピーライト・
セパレーターライン・テーマカラーが自動的に継承される。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# === 松尾研究所 カラーパレット ===
class C:
    bgWhite     = RGBColor(0xFF, 0xFF, 0xFF)
    bgLight     = RGBColor(0xF5, 0xF7, 0xFA)
    bgNavy      = RGBColor(0x00, 0x20, 0x60)
    textBlack   = RGBColor(0x00, 0x00, 0x00)
    textDark    = RGBColor(0x33, 0x33, 0x33)
    textMuted   = RGBColor(0x88, 0x88, 0x88)
    textWhite   = RGBColor(0xFF, 0xFF, 0xFF)
    tx2         = RGBColor(0x44, 0x54, 0x6A)
    accent1     = RGBColor(0x5B, 0x9B, 0xD5)
    accent2     = RGBColor(0xED, 0x7D, 0x31)
    accent3     = RGBColor(0xA5, 0xA5, 0xA5)
    accent4     = RGBColor(0xFF, 0xC0, 0x00)
    accent5     = RGBColor(0x44, 0x72, 0xC4)
    accent6     = RGBColor(0x70, 0xAD, 0x47)
    accent4Dark  = RGBColor(0xBF, 0x90, 0x00)
    accent4Light = RGBColor(0xFF, 0xD9, 0x66)
    accent5Dark  = RGBColor(0x33, 0x4F, 0x93)
    accent2Dark  = RGBColor(0xB2, 0x5C, 0x25)
    tx2Dark      = RGBColor(0x33, 0x3D, 0x4F)
    accentRed   = RGBColor(0xC8, 0x10, 0x2E)
    accentGreen = RGBColor(0x70, 0xAD, 0x47)
    gridLine    = RGBColor(0xD0, 0xD0, 0xD0)
    separator   = RGBColor(0x00, 0x00, 0x00)


# === レイアウト定数 ===
class L:
    slideW = Inches(13.33)
    slideH = Inches(7.5)
    mx = Inches(0.68)
    contentX = Inches(0.68)
    contentY = Inches(0.82)
    contentW = Inches(11.95)
    contentH = Inches(5.98)
    sourceX = Inches(0.68)
    sourceY = Inches(7.0)
    pageNumX = Inches(12.71)
    pageNumY = Inches(7.28)
    col2W = Inches(5.78)
    col2Gap = Inches(0.4)
    col2RightX = Inches(6.86)
    col3W = Inches(3.72)
    col3Gap = Inches(0.4)
    col3MidX = Inches(4.80)
    col3RightX = Inches(8.92)


# === レイアウトインデックス ===
LAYOUT_TITLE = 0
LAYOUT_CONTENT = 1
LAYOUT_FREE = 2
LAYOUT_SECTION = 3


# === ヘルパー関数 ===

def set_font(run, name, size, bold=False, color=None):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_text_box(slide, x, y, w, h, text, font_name="Hiragino Kaku Gothic Pro W3",
                 font_size=16, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, font_name, font_size, bold, color)
    return txBox


def add_header_bar(slide, x, y, w, text, fill_color=None, text_color=None, h=None):
    bar_h = h or Inches(0.40)
    fill = fill_color or C.bgNavy
    txt_color = text_color or C.textWhite
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, bar_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_font(run, "Hiragino Kaku Gothic Pro W3", 15, bold=True, color=txt_color)
    return shape


def add_number_badge(slide, x, y, num, size=None):
    s = size or Inches(0.50)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, s, s)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C.tx2
    shape.line.color.rgb = C.tx2Dark
    shape.line.width = Pt(1)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(num)
    set_font(run, "Hiragino Kaku Gothic Pro W3", 16, bold=True, color=C.textWhite)
    return shape


def clear_existing_slides(prs):
    """テンプレートの既存スライドを全削除

    Note: python-pptx にはスライド削除の公式APIがないため、
    内部XMLを直接操作している。python-pptx 1.0.2 で動作確認済み。
    """
    while len(prs.slides._sldIdLst) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


# === テンプレート関数 ===

def add_title_slide(prs, title, date=None, department=None, client_name=None, author=None):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    slide.placeholders[0].text = title
    if date and 1 in slide.placeholders:
        slide.placeholders[1].text = date
    if department and 10 in slide.placeholders:
        slide.placeholders[10].text = department
    if client_name and 11 in slide.placeholders:
        slide.placeholders[11].text = client_name
    if author and 12 in slide.placeholders:
        slide.placeholders[12].text = author
    return slide


def add_section_slide(prs, section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    slide.placeholders[0].text = section_title
    return slide


def add_free_content_slide(prs, title, source_text=None, page_num=None):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_FREE])
    slide.placeholders[0].text = title
    if page_num is not None and 12 in slide.placeholders:
        slide.placeholders[12].text = str(page_num)
    if source_text:
        txBox = slide.shapes.add_textbox(L.sourceX, L.sourceY, Inches(8), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"出典）{source_text}"
        set_font(run, "Hiragino Kaku Gothic Pro W3", 10.5, color=C.textDark)
    return slide


def add_styled_table(slide, headers, rows, x=None, y=None, w=None, col_widths=None,
                     header_font_size=12, body_font_size=10):
    tbl_x = x or L.contentX
    tbl_y = y or L.contentY
    tbl_w = w or L.contentW
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, tbl_x, tbl_y, tbl_w, Inches(0.1))
    table = table_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    # ヘッダー行
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                set_font(run, "Hiragino Kaku Gothic Pro W3", header_font_size,
                         bold=True, color=C.textWhite)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C.accent5Dark
    # データ行
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    set_font(run, "Hiragino Kaku Gothic Pro W3", body_font_size,
                             color=C.textBlack)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C.bgWhite if i % 2 == 0 else C.bgLight
    return table


# === プレゼンテーション構築 ===

def build_deck():
    import os
    import sys

    template_path = ".claude/skills/slides-maker/references/松尾研究所テンプレ_v2_DLしてお使いください.pptx"
    if not os.path.exists(template_path):
        print(f"エラー: テンプレートが見つかりません: {template_path}")
        print("プロジェクトルートから実行してください。")
        sys.exit(1)
    prs = Presentation(template_path)

    # 既存スライドを全削除
    clear_existing_slides(prs)

    # ──────────────────────────────────────────────
    # スライド1: タイトルスライド
    # ──────────────────────────────────────────────
    add_title_slide(prs,
        title="松尾研究所におけるAIツール導入ロードマップ",
        date="2026/03/02",
        department="松尾研究所 MLシステム開発チーム",
        author="岡田"
    )

    # ──────────────────────────────────────────────
    # スライド2: Agenda（自由配置レイアウト）
    # ──────────────────────────────────────────────
    slide2 = add_free_content_slide(prs, title="Agenda", page_num=2)
    agenda_items = [
        "設計思想 ─ 「点→線→面」の3段階アプローチ",
        "3段階ロードマップ ─ Stage A / B / C の役割分担",
        "カテゴリ別推奨マトリクス ─ 13カテゴリ×3ステージ",
        "規模帯別推奨ツール早見表",
        "松尾研の現在地 ─ 現状のツール利用状況とステージ判定",
        "想定される課題と乗り越え方",
        "展開順序とガバナンス ─ ネクストステップ",
    ]
    for i, item in enumerate(agenda_items):
        yBase = L.contentY + Inches(0.20) + Inches(i * 0.65)
        badgeSize = Inches(0.40)
        add_number_badge(slide2, L.contentX + Inches(0.3), yBase, i + 1, badgeSize)
        # テキストボックス（ボーダー付き）
        shape = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            L.contentX + Inches(0.88), yBase, Inches(10.0), badgeSize
        )
        shape.fill.background()
        shape.line.color.rgb = C.textBlack
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        set_font(run, "Hiragino Kaku Gothic Pro W3", 14, bold=True, color=C.textBlack)

    # ──────────────────────────────────────────────
    # スライド3: セクション区切り
    # ──────────────────────────────────────────────
    add_section_slide(prs, "1. 設計思想と3段階ロードマップ")

    # ──────────────────────────────────────────────
    # スライド4: 設計思想 ─ 3つの原則
    # ──────────────────────────────────────────────
    slide4 = add_free_content_slide(prs,
        title="「点→線→面」の段階展開が、AI導入の成功確率を最大化する",
        source_text="各社公式サイト・料金ページの公開情報に基づく整理（2026年2月時点）",
        page_num=4
    )

    dY = L.contentY + Inches(0.10)
    add_header_bar(slide4, L.contentX, dY, L.contentW, "設計思想 ─ 3つの原則")

    design_principles = [
        ("1", "保存先・フォーマットの統一",
         "Stage Aの段階からDrive/OneDrive等の保存先とOffice互換フォーマットを揃え、後段の移行コストを最小化する"),
        ("2", "統制（ガバナンス）の段階的強化",
         "AI導入の失敗要因は「モデル性能」よりも権限・共有・監査・データ利用の曖昧さに起因する。Stageが進むほどガバナンスを厚くする"),
        ("3", "専用ツールとスイートの使い分け",
         "会議要約はスイート内蔵で置換しやすい一方、CRM・会計・自動化等の基幹業務は専用ツールを残す領域を明確化する"),
    ]
    for i, (num, title, desc) in enumerate(design_principles):
        py = dY + Inches(0.52 + i * 0.62)
        add_number_badge(slide4, L.contentX + Inches(0.15), py, num, Inches(0.38))
        add_text_box(slide4, L.contentX + Inches(0.68), py, Inches(3.2), Inches(0.38),
                     title, "Hiragino Kaku Gothic Pro W6", 13, bold=True, color=C.textBlack)
        add_text_box(slide4, L.contentX + Inches(3.95), py, Inches(7.7), Inches(0.38),
                     desc, "Hiragino Kaku Gothic Pro W3", 11, color=C.textDark)

    # 3ステージ概要
    esY = dY + Inches(2.50)
    add_header_bar(slide4, L.contentX, esY, L.contentW,
                   "3段階ロードマップ ─ Stage A → B → C の役割分担")

    stages = [
        ("Stage A", "〜約30名", C.accent5, "点", "専用ツールで「点」の最適化",
         "痛点カテゴリを2〜4つ選び、専用ツールでROIを最短化。週次の削減時間と作業品質で効果測定"),
        ("Stage B", "30〜300名", C.accent6, "線", "Google Workspace＋Geminiで「線」の最適化",
         "日常アプリ（Gmail/Docs/Sheets/Meet等）にAIを埋め込み、コンテンツの保存先をDrive中心に統一"),
        ("Stage C", "300名以上", C.accent2, "面", "Microsoft 365＋Copilotで「面」の統制",
         "ID（Entra）・Graph・データ保護・監査を前提に全社展開。Power Automateで業務プロセスまで自動化"),
    ]
    for i, (label, subtitle, color, icon, desc, detail) in enumerate(stages):
        sx = L.contentX + Inches(i * (3.72 + 0.4))
        sy = esY + Inches(0.50)
        # ステージヘッダー
        shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, sy, L.col3W, Inches(0.42))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{label}（{subtitle}）"
        set_font(run, "Hiragino Kaku Gothic Pro W3", 13, bold=True, color=C.textWhite)

        # 「点/線/面」バッジ
        badge = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        sx + Inches(0.10), sy + Inches(0.52),
                                        Inches(0.50), Inches(0.50))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.fill.fore_color.brightness = 0.7
        badge.line.color.rgb = color
        badge.line.width = Pt(1)
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = icon
        set_font(run, "Hiragino Kaku Gothic Pro W6", 20, bold=True, color=color)

        # 説明
        add_text_box(slide4, sx + Inches(0.70), sy + Inches(0.52),
                     L.col3W - Inches(0.80), Inches(0.28),
                     desc, "Hiragino Kaku Gothic Pro W6", 11, bold=True, color=C.textBlack)
        add_text_box(slide4, sx + Inches(0.70), sy + Inches(0.80),
                     L.col3W - Inches(0.80), Inches(0.65),
                     detail, "Hiragino Kaku Gothic Pro W3", 10, color=C.textDark)

    # 移行トリガー注記
    add_text_box(slide4, L.contentX, esY + Inches(2.10), L.contentW, Inches(0.30),
                 "※ Stage B→Cは順番に移行するのではなく、①M&A・親会社標準 ②Microsoft中心のID/セキュリティ要件 ③Power Platform内製自動化の拡大 が揃ったタイミングで実行",
                 "Hiragino Kaku Gothic Pro W3", 9, color=C.textMuted)

    # ──────────────────────────────────────────────
    # スライド5: セクション区切り
    # ──────────────────────────────────────────────
    add_section_slide(prs, "2. カテゴリ別推奨マトリクス")

    # ──────────────────────────────────────────────
    # スライド6: カテゴリ別マトリクス Part 1
    # ──────────────────────────────────────────────
    slide6 = add_free_content_slide(prs,
        title="13カテゴリ×3ステージで最適ツールを配置する（前半7カテゴリ）",
        source_text="各社公式サイト・料金ページの公開情報（2026年2月時点）",
        page_num=6
    )
    add_styled_table(slide6,
        ["カテゴリ", "Stage A: 小規模（専用ツール）", "Stage B: Google＋Gemini", "Stage C: Microsoft＋Copilot"],
        [
            ["会議議事録", "Notta / Otter.ai / Fireflies.ai", "Meet「Take notes for me」\n→Docsに自動整理", "Teams Copilot\n要点・発言・アクション要約"],
            ["スライド生成", "Gamma / Beautiful.ai / Plus AI\nCanva（$20/人/月〜）", "Slides＋Gemini\n新規生成/要約/書き換え", "PowerPoint＋Copilot\nプロンプトからプレゼン生成"],
            ["文書ドラフト", "ChatGPT Business（¥3,900/月）\nClaude Team（$20/席/月）", "Gmail/Docs/Sheets\nサイドパネルでGemini支援", "Word/Excel/Outlook\nCopilot統合"],
            ["メール仕分け", "Shortwave（$24/席/月〜）\nFront（$25/席/月〜）", "Gmail＋Gemini\n受信箱検索・要点抽出", "Outlook＋Copilot\nメール文脈で要約・返信"],
            ["ナレッジ・検索", "Notion Business（¥3,150/人/月）\nGuru（$25/席/月〜）", "Drive中心＋NotebookLM\n管理者がAI機能を制御", "Graph典拠/コネクタ\nCopilot Chatで統合検索"],
            ["CRM・営業支援", "Pipedrive（$14/席/月〜）\nHubSpot / Close", "Gmail/Docs/Meet連携\nで提案書・メール生産性向上", "Sales in M365 Copilot\nOutlook/Teamsと統合"],
            ["カスタマーサポート", "Zendesk（$50/エージェント/月〜）\nIntercom Fin（$0.99/解決）", "Gmail返信下書き可能\n（チケット連動は専用継続）", "Service in M365 Copilot\nTeamsでケース要約"],
        ],
        y=L.contentY + Inches(0.08),
        col_widths=[1.7, 3.2, 3.4, 3.65],
        header_font_size=11, body_font_size=9
    )

    # ──────────────────────────────────────────────
    # スライド7: カテゴリ別マトリクス Part 2
    # ──────────────────────────────────────────────
    slide7 = add_free_content_slide(prs,
        title="13カテゴリ×3ステージで最適ツールを配置する（後半6カテゴリ）",
        source_text="各社公式サイト・料金ページの公開情報（2026年2月時点）",
        page_num=7
    )
    add_styled_table(slide7,
        ["カテゴリ", "Stage A: 小規模（専用ツール）", "Stage B: Google＋Gemini", "Stage C: Microsoft＋Copilot"],
        [
            ["開発者支援", "GitHub Copilot Business（$19/月）\nCursor（$20/月〜）", "Gemini Code Assist\n（$0〜$19〜$45）", "GitHub Copilot等併用\n（M365単体は開発支援弱）"],
            ["経理・財務", "freee / マネーフォワード\nQuickBooks", "周辺業務（帳票/申請）整備\n（会計自体は専用SaaS継続）", "Power Automate\n請求書処理・承認フロー自動化"],
            ["人事・採用", "Workable（$299/月〜）\nBreezy HR（$157/月〜）", "面接メモ/評価シート\nをDocsで標準化（ATS継続）", "Teams会議要約＋\nOutlook/SharePointで証跡統制"],
            ["法務・契約", "DocuSign（$10〜$25/月）\nJuro / LegalForce", "Docs＋Geminiでドラフト\nDriveで版管理", "Word＋Copilotでドラフト\neサイン/CLMは別製品併用"],
            ["分析・BI", "ThoughtSpot（$25/月〜）\nTableau", "Looker Studio＋Gemini\n自然言語クエリ", "Power BI＋Copilot\n（F2以上/P1以上キャパ要件）"],
            ["日程調整", "Calendly / Reclaim.ai（$12/月〜）\nClockwise（$6.75/月〜）", "Gmail候補日時提案\nメール起点の調整補助", "Outlook/Teams文脈\nで調整・要約"],
        ],
        y=L.contentY + Inches(0.08),
        col_widths=[1.7, 3.2, 3.4, 3.65],
        header_font_size=11, body_font_size=9
    )

    # ワークフロー自動化（13番目のカテゴリ）
    wfY = L.contentY + Inches(4.40)
    add_header_bar(slide7, L.contentX, wfY, L.contentW,
                   "＋ ワークフロー自動化・RPA（13番目のカテゴリ）", h=Inches(0.35))
    wf_items = [
        ("Stage A", "Zapier（無料〜）/ Make（$9〜）/ n8n / UiPath（$25/月〜）", C.accent5),
        ("Stage B", "AppSheet（$5/ユーザー/月〜）＋スプレッドシート連携", C.accent6),
        ("Stage C", "Power Automate Premium（¥2,248/ユーザー/月〜）＋RPA/プロセスマイニング", C.accent2),
    ]
    for i, (wf_label, wf_text, wf_color) in enumerate(wf_items):
        wx = L.contentX + Inches(i * (3.72 + 0.4))
        # カラーバー
        bar = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      wx, wfY + Inches(0.42), Inches(0.12), Inches(0.48))
        bar.fill.solid()
        bar.fill.fore_color.rgb = wf_color
        bar.line.fill.background()
        add_text_box(slide7, wx + Inches(0.20), wfY + Inches(0.42),
                     Inches(1.2), Inches(0.22),
                     wf_label, "Hiragino Kaku Gothic Pro W6", 10, bold=True, color=C.textBlack)
        add_text_box(slide7, wx + Inches(0.20), wfY + Inches(0.64),
                     L.col3W - Inches(0.35), Inches(0.28),
                     wf_text, "Hiragino Kaku Gothic Pro W3", 9, color=C.textDark)

    # ──────────────────────────────────────────────
    # スライド8: セクション区切り
    # ──────────────────────────────────────────────
    add_section_slide(prs, "3. 規模帯別推奨ツール早見表")

    # ──────────────────────────────────────────────
    # スライド9: 規模帯別早見表（3カラム）
    # ──────────────────────────────────────────────
    slide9 = add_free_content_slide(prs,
        title="各Stageで「何を主軸に」「何を専用ツールとして残すか」を明確化する",
        source_text="各社公式サイト・料金ページの公開情報に基づく整理（2026年2月時点）",
        page_num=9
    )

    stage_details = [
        {
            "title": "Stage A（〜約30名）", "subtitle": "「点」の最適化", "color": C.accent5,
            "focus": "痛点カテゴリを2〜4つ選定",
            "tools": ["会議: Notta / Otter / Fireflies", "資料: Gamma / Beautiful.ai",
                      "メール: Shortwave / Spark", "採用: Workable / Breezy", "CS: Zendesk / Intercom"],
            "keep": "CRM / サポート / 会計 / 採用 / 法務 / BI / 自動化（基幹に近い業務アプリ）"
        },
        {
            "title": "Stage B（30〜300名）", "subtitle": "Google Workspace＋Geminiで「線」の最適化",
            "color": C.accent6, "focus": "日常アプリの中でAIを回す",
            "tools": ["会議: Meet→Docs自動整理", "資料: Slides＋Gemini", "メール: Gmail＋Gemini",
                      "ナレッジ: Drive＋NotebookLM", "BI: Looker Studio＋Gemini"],
            "keep": "Zendesk等のCS / Pipedrive等のCRM / freee等の会計 / Zapier等の外部連携自動化"
        },
        {
            "title": "Stage C（300名以上）", "subtitle": "Microsoft 365＋Copilotで「面」の統制",
            "color": C.accent2, "focus": "ID・Graph・監査を前提に全社展開",
            "tools": ["会議: Teams＋Copilot", "資料: PowerPoint＋Copilot", "文書: Word＋Copilot",
                      "メール: Outlook＋Copilot", "自動化: Power Automate"],
            "keep": "GitHub Copilot等の開発支援 / CS・CRM・会計・法務の専門SaaS（連携して統制）"
        },
    ]

    for i, st in enumerate(stage_details):
        sx = L.contentX + Inches(i * (3.72 + 0.4))
        sy = L.contentY + Inches(0.10)

        # ステージヘッダー
        shape = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, sy, L.col3W, Inches(0.42))
        shape.fill.solid()
        shape.fill.fore_color.rgb = st["color"]
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = st["title"]
        set_font(run, "Hiragino Kaku Gothic Pro W3", 13, bold=True, color=C.textWhite)

        # サブタイトル
        add_text_box(slide9, sx + Inches(0.1), sy + Inches(0.50),
                     L.col3W - Inches(0.2), Inches(0.28),
                     st["subtitle"], "Hiragino Kaku Gothic Pro W6", 12,
                     bold=True, color=st["color"])

        # フォーカス
        focus_shape = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              sx, sy + Inches(0.85), L.col3W, Inches(0.30))
        focus_shape.fill.solid()
        focus_shape.fill.fore_color.rgb = st["color"]
        focus_shape.fill.fore_color.brightness = 0.85
        focus_shape.line.fill.background()
        add_text_box(slide9, sx + Inches(0.1), sy + Inches(0.85),
                     L.col3W - Inches(0.2), Inches(0.30),
                     f"主軸: {st['focus']}", "Hiragino Kaku Gothic Pro W3", 10,
                     bold=True, color=C.textBlack)

        # 推奨ツールリスト
        add_text_box(slide9, sx + Inches(0.1), sy + Inches(1.25),
                     L.col3W - Inches(0.2), Inches(0.22),
                     "推奨ツール:", "Hiragino Kaku Gothic Pro W6", 10,
                     bold=True, color=C.accent5Dark)

        tools_text = "\n".join(f"• {t}" for t in st["tools"])
        add_text_box(slide9, sx + Inches(0.1), sy + Inches(1.50),
                     L.col3W - Inches(0.2), Inches(1.70),
                     tools_text, "Hiragino Kaku Gothic Pro W3", 9, color=C.textBlack)

        # 専用ツール残存
        keep_header = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              sx, sy + Inches(3.30), L.col3W, Inches(0.25))
        keep_header.fill.solid()
        keep_header.fill.fore_color.rgb = C.accent4Light
        keep_header.fill.fore_color.brightness = 0.4
        keep_header.line.fill.background()
        add_text_box(slide9, sx + Inches(0.1), sy + Inches(3.30),
                     L.col3W - Inches(0.2), Inches(0.25),
                     "専用ツールを残す領域:", "Hiragino Kaku Gothic Pro W6", 9,
                     bold=True, color=C.textBlack)
        add_text_box(slide9, sx + Inches(0.1), sy + Inches(3.58),
                     L.col3W - Inches(0.2), Inches(0.75),
                     st["keep"], "Hiragino Kaku Gothic Pro W3", 9, color=C.textDark)

    # ──────────────────────────────────────────────
    # スライド10: セクション区切り — 松尾研の現在地 【新規】
    # ──────────────────────────────────────────────
    add_section_slide(prs, "4. 松尾研の現在地")

    # ──────────────────────────────────────────────
    # スライド11: 松尾研の現状ツール利用状況とステージ判定 【新規】
    # ──────────────────────────────────────────────
    slide11_new = add_free_content_slide(prs,
        title="松尾研のツール利用は「点」の段階 ─ Stage Bへの移行が次の一手",
        page_num=11
    )

    # 左カラム: 現状のツール利用状況
    lx = L.contentX
    ly = L.contentY + Inches(0.10)
    add_header_bar(slide11_new, lx, ly, L.col2W, "現状のツール利用状況")

    current_tools = [
        ("Notion", "ナレッジ・プロジェクト管理"),
        ("Slack", "コミュニケーション"),
        ("Google Workspace", "メール・カレンダー・Drive"),
        ("GitHub / Cursor / Claude Code", "開発支援"),
        ("ChatGPT / Claude", "文書ドラフト・分析"),
        ("Zoom / Google Meet", "会議"),
    ]
    add_styled_table(slide11_new,
        ["ツール", "用途"],
        current_tools,
        x=lx, y=ly + Inches(0.48), w=L.col2W,
        col_widths=[2.8, 2.98],
        header_font_size=11, body_font_size=10
    )

    # 右カラム: ステージ判定
    rx = L.col2RightX
    ry = ly
    add_header_bar(slide11_new, rx, ry, L.col2W, "ステージ判定")

    # 現在
    sj_y = ry + Inches(0.58)
    badge_current = slide11_new.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                  rx + Inches(0.10), sj_y,
                                                  Inches(1.0), Inches(0.35))
    badge_current.fill.solid()
    badge_current.fill.fore_color.rgb = C.accent5
    badge_current.line.fill.background()
    tf = badge_current.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "現在"
    set_font(run, "Hiragino Kaku Gothic Pro W3", 12, bold=True, color=C.textWhite)

    add_text_box(slide11_new, rx + Inches(1.25), sj_y, Inches(4.4), Inches(0.35),
                 "Stage A寄り（専用ツールの「点」利用が中心）",
                 "Hiragino Kaku Gothic Pro W3", 12, bold=True, color=C.textBlack)

    # 課題
    sj_y2 = sj_y + Inches(0.55)
    badge_issue = slide11_new.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                rx + Inches(0.10), sj_y2,
                                                Inches(1.0), Inches(0.35))
    badge_issue.fill.solid()
    badge_issue.fill.fore_color.rgb = C.accent2
    badge_issue.line.fill.background()
    tf = badge_issue.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "課題"
    set_font(run, "Hiragino Kaku Gothic Pro W3", 12, bold=True, color=C.textWhite)

    add_text_box(slide11_new, rx + Inches(1.25), sj_y2, Inches(4.4), Inches(0.70),
                 "ツール間の連携が弱く、成果物の保存先がバラバラ\n（Notion / Drive / ローカル が混在）",
                 "Hiragino Kaku Gothic Pro W3", 11, color=C.textDark)

    # 目指す姿
    sj_y3 = sj_y2 + Inches(0.85)
    badge_goal = slide11_new.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                               rx + Inches(0.10), sj_y3,
                                               Inches(1.0), Inches(0.35))
    badge_goal.fill.solid()
    badge_goal.fill.fore_color.rgb = C.accent6
    badge_goal.line.fill.background()
    tf = badge_goal.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "目指す姿"
    set_font(run, "Hiragino Kaku Gothic Pro W3", 11, bold=True, color=C.textWhite)

    add_text_box(slide11_new, rx + Inches(1.25), sj_y3, Inches(4.4), Inches(0.70),
                 "Stage Bへの移行\nGoogle Workspace + Gemini中心の「線」の最適化",
                 "Hiragino Kaku Gothic Pro W6", 12, bold=True, color=C.accent6)

    # 補足: 移行イメージ矢印
    arrow_y = sj_y3 + Inches(1.0)
    arrow_shape = slide11_new.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                rx + Inches(0.10), arrow_y,
                                                L.col2W - Inches(0.20), Inches(0.40))
    arrow_shape.fill.solid()
    arrow_shape.fill.fore_color.rgb = C.bgLight
    arrow_shape.line.fill.background()
    tf = arrow_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Stage A（点）  →  Stage B（線）  →  Stage C（面）"
    set_font(run, "Hiragino Kaku Gothic Pro W3", 11, bold=True, color=C.tx2)

    # ──────────────────────────────────────────────
    # スライド12: セクション区切り — 想定される課題と乗り越え方 【新規】
    # ──────────────────────────────────────────────
    add_section_slide(prs, "5. 想定される課題と乗り越え方")

    # ──────────────────────────────────────────────
    # スライド13: 課題5項目テーブル + 乗り越え方 【新規】
    # ──────────────────────────────────────────────
    slide13 = add_free_content_slide(prs,
        title="導入の壁を事前に想定し、段階的に乗り越える",
        page_num=13
    )

    add_styled_table(slide13,
        ["課題", "具体的な壁", "乗り越え方"],
        [
            ["既存ワークフロー\nとの摩擦",
             "既存ツール（グループウェア・ファイル\nサーバー等）の運用が定着しており、\nDrive/SharePointへの移行は大きな変更コスト",
             "既存ツールを残しつつクラウド\nストレージを「保存先」として併用。\n段階的にDrive/OneDrive中心へ移行"],
            ["コスト・予算\nの確保",
             "Gemini/Copilotのライセンス費\n（$20〜30/人/月）、PoC期間の\n既存ツールとの二重運用コスト",
             "少人数（5〜10名）のパイロット\nチームで効果検証し、\nROIを定量化してから全社拡大"],
            ["セキュリティ・\nデータ統制",
             "社内機密情報・顧客データのAI送信\nに関する情報セキュリティポリシー\nとの整合性、データ漏洩リスク",
             "ツールごとの学習/保持ポリシーを\n整理し、利用可能な範囲を明文化。\n機密度に応じた3段階ルールを設定"],
            ["定着・教育\nコスト",
             "現場社員は多忙で新ツール習得の\n時間が限られる。「使う人」と\n「使わない人」の二極化リスク",
             "「会議AI」「メールAI」など日常\n業務に直結する領域から開始し、\n初期成功体験を社内で共有"],
            ["効果測定\nの難しさ",
             "ホワイトカラー業務の生産性は\n定量化しにくい（企画の質、\n意思決定の速度等）",
             "定量指標（会議議事録作成時間、\nメール処理時間）と定性指標\n（満足度アンケート）を併用"],
        ],
        y=L.contentY + Inches(0.08),
        col_widths=[1.8, 4.2, 5.95],
        header_font_size=11, body_font_size=9
    )

    # ──────────────────────────────────────────────
    # スライド14: セクション区切り — 展開順序とガバナンス
    # ──────────────────────────────────────────────
    add_section_slide(prs, "6. 展開順序とガバナンス")

    # ──────────────────────────────────────────────
    # スライド15: 展開順序 + ガバナンス
    # ──────────────────────────────────────────────
    slide11 = add_free_content_slide(prs,
        title="「機能から」ではなく「統制から」始めることが、展開の失敗確率を下げる",
        source_text="Google Workspace/Microsoft 365 公式ドキュメント・料金ページに基づく整理",
        page_num=15
    )

    eoY = L.contentY + Inches(0.10)
    add_header_bar(slide11, L.contentX, eoY, L.contentW,
                   "推奨展開順序 ─ 定着しやすい領域から段階的に拡大", h=Inches(0.38))

    deploy_order = [
        ("1", "会議（要約/議事録）", "成果物が残りやすく、定着の初速が出る"),
        ("2", "メール要約", "毎日の業務に直結し、効果実感が早い"),
        ("3", "文書/スライドのテンプレ運用", "テンプレ統一・版管理のニーズが自然に発生"),
        ("4", "ナレッジ検索", "「あの資料どこ？」の属人化を解消"),
        ("5", "顧客接点（CRM/サポート）", "営業工数削減・CS自動化で投資対効果が明確"),
        ("6", "自動化（RPA/ワークフロー）とBI", "業務プロセス全体の最適化で全社横断展開"),
    ]
    for i, (num, title, reason) in enumerate(deploy_order):
        dy = eoY + Inches(0.48 + i * 0.48)
        add_number_badge(slide11, L.contentX + Inches(0.15), dy + Inches(0.02), num, Inches(0.32))
        add_text_box(slide11, L.contentX + Inches(0.60), dy, Inches(3.5), Inches(0.36),
                     title, "Hiragino Kaku Gothic Pro W6", 12, bold=True, color=C.textBlack)
        add_text_box(slide11, L.contentX + Inches(4.20), dy, Inches(7.5), Inches(0.36),
                     reason, "Hiragino Kaku Gothic Pro W3", 11, color=C.textDark)
        # 矢印
        if i < len(deploy_order) - 1:
            add_text_box(slide11, L.contentX + Inches(0.20), dy + Inches(0.36),
                         Inches(0.22), Inches(0.12),
                         "▼", "Calibri", 7, color=C.accent3, align=PP_ALIGN.CENTER)

    # ガバナンス標準セット
    gvY = eoY + Inches(3.45)
    add_header_bar(slide11, L.contentX, gvY, L.contentW,
                   "ガバナンス標準セット ─ 最低限準備すべき3項目", h=Inches(0.38))

    gov_items = [
        ("機密データの取り扱い",
         "ツールごとの学習/保持ポリシーを確認し（例: ChatGPT Business/Claude Teamは既定で学習しない）、利用範囲（社外秘/個人情報/契約書等）を定義"),
        ("自動化のコネクタ/DLP設計",
         "M365標準コネクタのみの制約とプレミアムコネクタ要件を明確化し、コネクタ種別とDLPを設計してから拡大"),
        ("外部データ共有の同意フロー",
         "CS/営業など外部データ共有が発生する領域は、管理者承認の同意フローを必須化（Sales統合等）"),
    ]
    for i, (g_title, g_desc) in enumerate(gov_items):
        gy = gvY + Inches(0.46 + i * 0.58)
        badge = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         L.contentX + Inches(0.15), gy,
                                         Inches(0.32), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = C.accentRed
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        set_font(run, "Hiragino Kaku Gothic Pro W3", 12, bold=True, color=C.textWhite)

        add_text_box(slide11, L.contentX + Inches(0.60), gy, Inches(2.8), Inches(0.35),
                     g_title, "Hiragino Kaku Gothic Pro W6", 11, bold=True, color=C.textBlack)
        add_text_box(slide11, L.contentX + Inches(3.50), gy, Inches(8.2), Inches(0.50),
                     g_desc, "Hiragino Kaku Gothic Pro W3", 10, color=C.textDark)

    # ──────────────────────────────────────────────
    # スライド16: 推奨ネクストステップ（まとめ）
    # ──────────────────────────────────────────────
    slide12 = add_free_content_slide(prs,
        title="MLシステム開発チームをパイロットに、会議AI・メールAIからQuick Winを創出する",
        source_text="松尾研究所 MLシステム開発チーム",
        page_num=16
    )

    # サマリーKPI
    kpiY = L.contentY + Inches(0.15)
    kpis = [
        ("13", "対象カテゴリ"),
        ("3段階", "Stage A→B→C"),
        ("6〜18ヶ月", "段階展開期間"),
        ("点→線→面", "最適化アプローチ"),
    ]
    kpi_col_w = Inches(11.95 / 4)
    for i, (value, label) in enumerate(kpis):
        kpiX = L.contentX + Inches(i * (11.95 / 4))
        add_text_box(slide12, kpiX, kpiY, kpi_col_w, Inches(0.60),
                     value, "Calibri", 42, bold=True, color=C.accent5, align=PP_ALIGN.CENTER)
        add_text_box(slide12, kpiX, kpiY + Inches(0.62), kpi_col_w, Inches(0.22),
                     label, "Hiragino Kaku Gothic Pro W3", 10, color=C.textDark,
                     align=PP_ALIGN.CENTER)
        # セパレーター
        if i < len(kpis) - 1:
            sep = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           kpiX + kpi_col_w, kpiY + Inches(0.08),
                                           Pt(0.5), Inches(0.68))
            sep.fill.solid()
            sep.fill.fore_color.rgb = C.gridLine
            sep.line.fill.background()

    # 推奨ネクストステップ
    nsY = kpiY + Inches(1.10)
    add_header_bar(slide12, L.contentX, nsY, L.contentW, "推奨ネクストステップ")

    next_steps = [
        ("1", "松尾研の痛点カテゴリ選定",
         "会議AI・メールAI・ナレッジ検索を候補に、最も工数が大きい2〜3カテゴリを特定",
         "1週間以内"),
        ("2", "パイロットチーム（MLシステム開発チーム）でPoC開始",
         "選定ツールでPoCを実施。「週次の削減時間」と「作業品質」で効果測定し、ROIを定量化",
         "2ヶ月以内"),
        ("3", "保存先の統一ルール策定",
         "Notion＋Driveの併用ルールを整備し、成果物の保存先とフォーマットを段階的に統一",
         "1ヶ月以内"),
        ("4", "セキュリティ・学習ポリシーの明文化",
         "研究データ・論文草稿のAI送信ルールを整理し、大学IT規定との整合性を確認",
         "2ヶ月以内"),
        ("5", "Stage B移行計画の策定",
         "Google Workspace＋Gemini中心の「線」の最適化に向け、移行スケジュール・予算・教育プランを策定",
         "3ヶ月以内"),
    ]
    for i, (num, step_title, step_desc, timing) in enumerate(next_steps):
        sy = nsY + Inches(0.52 + i * 0.72)
        add_number_badge(slide12, L.contentX + Inches(0.15), sy, num, Inches(0.38))
        add_text_box(slide12, L.contentX + Inches(0.68), sy, Inches(2.8), Inches(0.38),
                     step_title, "Hiragino Kaku Gothic Pro W6", 12,
                     bold=True, color=C.textBlack)
        add_text_box(slide12, L.contentX + Inches(3.55), sy, Inches(6.6), Inches(0.50),
                     step_desc, "Hiragino Kaku Gothic Pro W3", 10, color=C.textDark)
        # タイミングバッジ
        timing_shape = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                L.contentX + Inches(10.35), sy + Inches(0.04),
                                                Inches(1.40), Inches(0.30))
        timing_shape.fill.solid()
        timing_shape.fill.fore_color.rgb = C.accent4Light
        timing_shape.line.color.rgb = C.accent4Dark
        timing_shape.line.width = Pt(0.5)
        tf = timing_shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = timing
        set_font(run, "Hiragino Kaku Gothic Pro W3", 9, bold=True, color=C.textBlack)

    # ──────────────────────────────────────────────
    # 保存
    # ──────────────────────────────────────────────
    output_path = "ai-tool-roadmap.pptx"
    prs.save(output_path)
    print(f"プレゼンテーションを保存しました: {output_path}")


if __name__ == "__main__":
    build_deck()
